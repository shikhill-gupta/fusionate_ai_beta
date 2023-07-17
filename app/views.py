"""
Definition of views.
"""



#DJANGO RELATED IMPORTS
from django import *
from django.shortcuts import *
from django.http import *
from .forms import *


#LANGCHAIN RELATED IMPORTS
from langchain import *
from langchain.document_loaders import PDFMinerLoader, Docx2txtLoader, CSVLoader
from langchain.chains import RetrievalQA
from langchain.llms import OpenAI
from langchain.indexes import VectorstoreIndexCreator
from bardapi import Bard
from langchain.vectorstores import Chroma
from langchain.embeddings import OpenAIEmbeddings
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.document_loaders import TextLoader
from langchain.document_loaders import DirectoryLoader, CSVLoader
from langchain.agents import create_csv_agent


#PYTHON RELATED IMPORTS
from datetime import datetime
import json
import os
import shutil

#DATA ANALYSIS IMPORTS
import pandas as pd

os.environ["OPENAI_API_KEY"] = "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"
os.environ["HUGGINGFACEHUB_API_TOKEN"] = "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"


#TOKENS


from rest_framework.generics import ListAPIView
from .models import File
from .serializers import FileSerializer
from pptx import Presentation


class FileListAPIView(ListAPIView):
    queryset = File.objects.all()
    serializer_class = FileSerializer
    lookup_field = 'user__pk'


class AITools:

    def __init__(self):
        pass


    def create_presentation(self, text_list):
        # Create a presentation object
        presentation = Presentation()

        # Iterate through the list of text and create a slide for each item
        for text in text_list:
            # Add a slide with a title and content layout
            slide_layout = presentation.slide_layouts[1]
            slide = presentation.slides.add_slide(slide_layout)

            # Set the title of the slide
            title = slide.shapes.title
            title.text = "Slide Title"

            # Set the content of the slide
            content = slide.placeholders[1]
            content.text = text

        # Save the presentation
        presentation.save("app/static/sample.pptx")




    def splitCharacters(self, my_string):
        #my_string = "Sed ut perspiciatis unde omnis iste natus error sit voluptatem accusantium doloremque laudantium, totam rem aperiam, eaque ipsa quae ab illo inventore veritatis et quasi architecto beatae vitae dicta sunt explicabo. Nemo enim ipsam voluptatem quia voluptas sit aspernatur aut odit aut fugit, sed quia consequuntur magni dolores eos qui ratione voluptatem sequi nesciunt. Neque porro quisquam est, qui dolorem ipsum quia dolor sit amet, consectetur, adipisci velit, sed quia non numquam eius modi tempora incidunt ut labore et dolore magnam aliquam quaerat voluptatem. Ut enim ad minima veniam, quis nostrum exercitationem ullam corporis suscipit laboriosam, nisi ut aliquid ex ea commodi consequatur? Quis autem vel eum iure reprehenderit qui in ea voluptate velit esse quam nihil molestiae consequatur, vel illum qui dolorem eum fugiat quo voluptas nulla pariatur?"
        splitted = my_string.split(" ")
        empty_string = ""
        empty_list = []
        counter = 0

        for x in splitted:
            empty_string = empty_string + str(x) + " "
            
            if counter == 15:
                print(empty_string, "THE EMPTY STRING IS 5 NOW RESETTING")
                empty_list.append(empty_string)
                empty_string = ""
                counter = 0
            else:
                print(empty_string, "IS THE EMPTY STRING")
            print("ADDING THE COUNTER THE COUNTER IS EQUAL TO", counter)
            counter = counter + 1
        return empty_list


class ChatWithData:

    def analyzeCSV(self, request, question, file_path):
        print("Analyze CSV Invoked.")

        try:
            print(file_path, "IS THE FILE PATH")

            question = request.POST.get('question')
            file_name = request.POST.get('file_name')
            user_folder = str(request.user)

            file_path = os.path.join('static', user_folder, file_name)
            file_dir = os.path.dirname(file_path)

            file_extension = os.path.splitext(file_name)[1]


            final_dir = str(file_dir) + "/" + os.path.splitext(file_name)[0] + ".csv"
            data = pd.read_csv(file_path)

            agent = create_csv_agent(OpenAI(temperature=0.9), 
                         str(final_dir), 
                         verbose=True)

            response = agent.run(question)
            
            return str(response)

        except Exception as e:
            print(e)
            return str(e)

    def analyzePDF(self, request, question, file_path):
        print("Analyze PDF Invoked.")

        try:
            loader = PDFMinerLoader(file_path)
            documents = loader.load()
            text_splitter = RecursiveCharacterTextSplitter(chunk_size=2048, chunk_overlap=0)
            texts = text_splitter.split_documents(documents)
            persist_directory = 'static/' + str(request.user) + "/db"
            embedding = OpenAIEmbeddings()

            vectordb = Chroma.from_documents(documents=texts, 
                                             embedding=embedding,
                                             persist_directory=persist_directory)

            retriever = vectordb.as_retriever(search_kwargs={"k": 3, "include_query": False})
            retriever.search_type
            retriever.search_kwargs
            chain = RetrievalQA.from_chain_type(llm=OpenAI(temperature=0.9), chain_type="stuff", retriever=retriever, input_key="question", return_source_documents=True)
            query = str(question)
            response = chain({"question": query})
            return str(response['result'])

        except Exception as e:
            print(e)
            return str(e)

    def analyzeDOCX(self, request, question, file_path):
        print("Analyze DOCX Invoked.")

        try:
            print(file_path, "IS THE FILE PATH")
            loader = Docx2txtLoader(file_path)
            documents = loader.load()
            text_splitter = RecursiveCharacterTextSplitter(chunk_size=2048, chunk_overlap=0)
            texts = text_splitter.split_documents(documents)
            persist_directory = 'static/' + str(request.user) + "/db"
            embedding = OpenAIEmbeddings()

            vectordb = Chroma.from_documents(documents=texts, 
                                             embedding=embedding,
                                             persist_directory=persist_directory)

            retriever = vectordb.as_retriever(search_kwargs={"k": 10, "include_query": False})
            retriever.search_type
            retriever.search_kwargs
            chain = RetrievalQA.from_chain_type(llm=OpenAI(temperature=0.9), chain_type="stuff", retriever=retriever, input_key="question", return_source_documents=True)
            query = str(question)
            response = chain({"question": query})

            A = AITools()
            answered = str(response['result'])
            print(answered, "IS THE ANSWER", type(answered), "IS THE TYPE")
            returned_list = A.splitCharacters(answered)
            print(returned_list, type(returned_list), "IS THE TYPE", "AND ITS VALUE")
            A.create_presentation(returned_list)

            giving_response = response['result']
            giving_response = giving_response  + "<a href='static/sample.pptx' download='sample.pptx'><button class='btn btn-warning'>Download the Presentation</button></a>"
            return giving_response

        except Exception as e:
            print(e)
            return str(e)

    def analyzeXLSX(self, request, question, file_path):
        print("Analyze Excel Invoked.")

        try:
            print(file_path, "IS THE FILE PATH")

            question = request.POST.get('question')
            file_name = request.POST.get('file_name')
            user_folder = str(request.user)

            file_path = os.path.join('static', user_folder, file_name)
            file_dir = os.path.dirname(file_path)

            file_extension = os.path.splitext(file_name)[1]


            final_dir = str(file_dir) + "/" + os.path.splitext(file_name)[0] + ".csv"
            data = pd.read_excel(file_path, engine='openpyxl')
            data.to_csv(final_dir, index=False)
            agent = create_csv_agent(OpenAI(temperature=0.9), 
                         str(final_dir), 
                         verbose=True)

            response = agent.run(question)
            
            return str(response)

        except Exception as e:
            print(e)
            return str(e)
class AISection:

    def __init__(self):
        pass


    def giveResponse(self, response_message):

        list_of_response = [response_message]
        json_response = {
            "message": list_of_response
        }

        json_string = json.dumps(json_response)
        response = HttpResponse(json_string, content_type='application/json')
        return response


    def useBard(self, message, api_key = 'XAgmCMSjzlW6gfVCdohc3Su_SHJj-OJ8nHgLzWRT_RJzBDg76z6jgr83kmx4D0zoXCsd6A.'):
        try:
            os.environ['_BARD_API_KEY'] = api_key
            my_question = str(message)
            return Bard().get_answer(my_question)['content']
        except Exception as e:
            print(e)
            return str(e)



    def useChatwithdataWithVectorChroma(self, request):

        question = request.POST.get('question')
        file_name = request.POST.get('file_name')
        user_folder = str(request.user)

        file_path = os.path.join('static', user_folder, file_name)
        file_dir = os.path.dirname(file_path)
        file_extension = os.path.splitext(file_name)[1]

        CD = ChatWithData()

        if file_extension == '.csv':
            return CD.analyzeCSV(request, question, file_path)
        if file_extension == '.pdf':
            return CD.analyzePDF(request, question, file_path)            
        if file_extension == '.docx':
            return CD.analyzeDOCX(request, question, file_path)        
        if file_extension == '.xlsx':
            return CD.analyzeXLSX(request, question, file_path)   

    def useHuggingface(self, message, define_model = "google/flan-t5-base"):

        try:

            llm = HuggingFaceHub(repo_id = define_model, model_kwargs={"temperature":0.9})
            template = """Question: {question} Let's think step by step. Answer: """
            prompt = PromptTemplate(template=template, input_variables=["question"])
            llm_chain = LLMChain(prompt=prompt, llm=llm)
            answer = llm_chain.run(message)
            print(answer, "IS THE ANSWER")
            return answer
        except Exception as e:
            print(e)
            return str(e)


def askQuestion(request):


    if request.method == 'POST':
        if request.POST.get('file_state') == 'yes':
            MyAI = AISection()
            response = MyAI.useChatwithdataWithVectorChroma(request)

            
            try:
                persist_directory = 'static/' + str(request.user) + "/db"
                shutil.rmtree(persist_directory)
            except Exception as e:
                pass

            return MyAI.giveResponse(response)


        else:
            print(request.POST.get('question'), "IS THE QUESTION")
            MyAI = AISection()
            response = MyAI.useBard(request.POST.get('question'))
            import re

            matches = re.findall(r"```([\s\S]*?)```", response)
            values = [match.strip() for match in matches]
            print(len(values))
            if len(values) > 0:


                return MyAI.giveResponse(values)
            else:
                return MyAI.giveResponse(response)




def handle_file_upload(request):
    print("IS TRIGGERED")
    if request.method == 'POST':
        form = FileUploadForm(request.POST, request.FILES)
        if form.is_valid():
            file = form.save(commit=False)
            file.file = request.FILES['file']
            file.save()
            print("SAVED")
        else:
            print("IS INVALID")
    else:
        form = FileUploadForm()
        print("NOT A POST")
    return redirect('home')

def home(request):


    from .forms import FileUploadForm
    """Renders the home page."""
    assert isinstance(request, HttpRequest)
    return render(
        request,
        'app/index.html',
        {
            'upload_file' : FileUploadForm(),
            'title':'Home Page',
            'year':datetime.now().year,
        }
    )

def contact(request):
    """Renders the contact page."""
    assert isinstance(request, HttpRequest)
    return render(
        request,
        'app/contact.html',
        {
            'title':'Contact',
            'message':'Your contact page.',
            'year':datetime.now().year,
        }
    )

def about(request):
    """Renders the about page."""
    assert isinstance(request, HttpRequest)
    return render(
        request,
        'app/about.html',
        {
            'title':'About',
            'message':'Your application description page.',
            'year':datetime.now().year,
        }
    )


def register(request):
    from .forms import RegistrationForm
    list_of_request = []
    for x in request.POST:
        list_of_request.append(str(request.POST.get(x)))

    if request.method == 'POST':

        form = RegistrationForm(request.POST)
        if form.is_valid():
            form.save()
            print("Valid")
            return shortcuts.redirect('login')
        else:
            print("Invalid")
    else:
        form = RegistrationForm()

    return shortcuts.redirect('login')